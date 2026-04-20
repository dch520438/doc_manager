"""
文档管家 - 核心模块
包含：数据库管理、规则引擎、文件整理器、内容提取器
"""

import os
import re
import sys
import platform
import shutil
import zipfile
import sqlite3
import json
import time
import logging
import threading
import hashlib
import subprocess
from datetime import datetime
from pathlib import Path
from typing import List, Dict, Optional, Tuple
from dataclasses import dataclass, asdict
from enum import Enum

# ==================== 平台检测 ====================
IS_WINDOWS = platform.system() == 'Windows'
IS_LINUX = platform.system() == 'Linux'
IS_MACOS = platform.system() == 'Darwin'

# ==================== 日志配置 ====================
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s [%(levelname)s] %(message)s',
    datefmt='%H:%M:%S'
)
logger = logging.getLogger("DocManager")


# ==================== 数据模型 ====================
@dataclass
class Rule:
    """分类规则"""
    id: int = 0
    name: str = ""               # 规则名称
    contain_keywords: str = ""    # 包含关键词（逗号分隔，全部满足）
    exclude_keywords: str = ""    # 排除关键词（逗号分隔，任一满足则排除）
    target_folder: str = ""       # 目标文件夹路径
    enabled: bool = True
    priority: int = 0             # 优先级，数字越小优先级越高
    created_at: str = ""

    def to_dict(self):
        return asdict(self)

    @classmethod
    def from_dict(cls, d):
        return cls(**d)


@dataclass
class Document:
    """文档记录"""
    id: int = 0
    filename: str = ""
    filepath: str = ""
    target_folder: str = ""
    rule_name: str = ""
    file_size: int = 0
    file_type: str = ""
    content_preview: str = ""
    file_hash: str = ""
    organized_at: str = ""
    original_path: str = ""


# ==================== 支持的文件格式 ====================
SUPPORTED_EXTENSIONS = {
    '.doc', '.docx', '.wps',          # Word / WPS 文字
    '.xls', '.xlsx', '.csv',          # Excel / 表格
    '.et', '.ett',                    # WPS 表格
    '.pdf',                            # PDF
    '.ppt', '.pptx', '.dps',          # PPT / WPS 演示
    '.txt', '.md', '.log',            # 纯文本
    '.rtf',                            # 富文本
    '.xml', '.json', '.html', '.htm', # 标记语言
}

EXTENSION_TYPE_MAP = {
    '.doc': 'Word', '.docx': 'Word', '.wps': 'WPS',
    '.xls': 'Excel', '.xlsx': 'Excel', '.csv': 'CSV',
    '.et': 'WPS表格', '.ett': 'WPS表格',
    '.pdf': 'PDF',
    '.ppt': 'PPT', '.pptx': 'PPT', '.dps': 'WPS演示',
    '.txt': '文本', '.md': 'Markdown', '.log': '日志',
    '.rtf': 'RTF',
    '.xml': 'XML', '.json': 'JSON', '.html': 'HTML', '.htm': 'HTML',
}

# 可提取文本内容的扩展名（用于搜索）
TEXT_EXTRACTABLE_EXTENSIONS = set(SUPPORTED_EXTENSIONS)


# ==================== 数据库管理 ====================
class Database:
    """SQLite数据库管理"""

    def __init__(self, db_path: str = None):
        if db_path is None:
            db_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "doc_manager.db")
        self.db_path = db_path
        self._lock = threading.Lock()
        self._conn = None
        self._init_db()

    def _get_conn(self):
        if self._conn is not None:
            return self._conn
        conn = sqlite3.connect(self.db_path, check_same_thread=False)
        conn.row_factory = sqlite3.Row
        if self.db_path != ':memory:':
            conn.execute("PRAGMA journal_mode=WAL")
        self._conn = conn
        return conn

    def _init_db(self):
        with self._lock:
            conn = self._get_conn()
            try:
                conn.executescript("""
                    CREATE TABLE IF NOT EXISTS rules (
                        id INTEGER PRIMARY KEY AUTOINCREMENT,
                        name TEXT NOT NULL,
                        contain_keywords TEXT DEFAULT '',
                        exclude_keywords TEXT DEFAULT '',
                        target_folder TEXT NOT NULL,
                        enabled INTEGER DEFAULT 1,
                        priority INTEGER DEFAULT 0,
                        created_at TEXT DEFAULT (datetime('now','localtime'))
                    );

                    CREATE TABLE IF NOT EXISTS documents (
                        id INTEGER PRIMARY KEY AUTOINCREMENT,
                        filename TEXT NOT NULL,
                        filepath TEXT NOT NULL,
                        target_folder TEXT NOT NULL,
                        rule_name TEXT DEFAULT '',
                        file_size INTEGER DEFAULT 0,
                        file_type TEXT DEFAULT '',
                        content_preview TEXT DEFAULT '',
                        file_hash TEXT DEFAULT '',
                        organized_at TEXT DEFAULT (datetime('now','localtime')),
                        original_path TEXT DEFAULT ''
                    );

                    CREATE TABLE IF NOT EXISTS settings (
                        key TEXT PRIMARY KEY,
                        value TEXT NOT NULL
                    );

                    CREATE INDEX IF NOT EXISTS idx_docs_filename ON documents(filename);
                    CREATE INDEX IF NOT EXISTS idx_docs_content ON documents(content_preview);
                    CREATE INDEX IF NOT EXISTS idx_docs_folder ON documents(target_folder);
                    CREATE INDEX IF NOT EXISTS idx_docs_type ON documents(file_type);
                """)
            finally:
                pass  # Don't close - keep connection alive

    # ---- 规则 CRUD ----
    def add_rule(self, rule: Rule) -> int:
        with self._lock:
            conn = self._get_conn()
            try:
                cursor = conn.execute(
                    "INSERT INTO rules (name, contain_keywords, exclude_keywords, target_folder, enabled, priority) VALUES (?,?,?,?,?,?)",
                    (rule.name, rule.contain_keywords, rule.exclude_keywords, rule.target_folder, int(rule.enabled), rule.priority)
                )
                conn.commit()
                return cursor.lastrowid
            finally:
                pass  # persistent connection

    def update_rule(self, rule: Rule):
        with self._lock:
            conn = self._get_conn()
            try:
                conn.execute(
                    "UPDATE rules SET name=?, contain_keywords=?, exclude_keywords=?, target_folder=?, enabled=?, priority=? WHERE id=?",
                    (rule.name, rule.contain_keywords, rule.exclude_keywords, rule.target_folder, int(rule.enabled), rule.priority, rule.id)
                )
                conn.commit()
            finally:
                pass  # persistent connection

    def delete_rule(self, rule_id: int):
        with self._lock:
            conn = self._get_conn()
            try:
                conn.execute("DELETE FROM rules WHERE id=?", (rule_id,))
                conn.commit()
            finally:
                pass  # persistent connection

    def get_all_rules(self) -> List[Rule]:
        conn = self._get_conn()
        try:
            rows = conn.execute("SELECT * FROM rules ORDER BY priority ASC, id ASC").fetchall()
            return [Rule.from_dict(dict(r)) for r in rows]
        finally:
            pass

    def get_enabled_rules(self) -> List[Rule]:
        conn = self._get_conn()
        try:
            rows = conn.execute("SELECT * FROM rules WHERE enabled=1 ORDER BY priority ASC, id ASC").fetchall()
            return [Rule.from_dict(dict(r)) for r in rows]
        finally:
            pass

    # ---- 文档记录 ----
    def add_document(self, doc: Document) -> int:
        with self._lock:
            conn = self._get_conn()
            try:
                cursor = conn.execute(
                    """INSERT INTO documents (filename, filepath, target_folder, rule_name, file_size, file_type, content_preview, file_hash, organized_at, original_path)
                       VALUES (?,?,?,?,?,?,?,?,?,?)""",
                    (doc.filename, doc.filepath, doc.target_folder, doc.rule_name, doc.file_size, doc.file_type, doc.content_preview, doc.file_hash, doc.organized_at, doc.original_path)
                )
                conn.commit()
                return cursor.lastrowid
            finally:
                pass  # persistent connection

    def delete_document(self, doc_id: int):
        with self._lock:
            conn = self._get_conn()
            try:
                conn.execute("DELETE FROM documents WHERE id=?", (doc_id,))
                conn.commit()
            finally:
                pass  # persistent connection

    def delete_documents_by_folder(self, folder: str):
        with self._lock:
            conn = self._get_conn()
            try:
                conn.execute("DELETE FROM documents WHERE target_folder=?", (folder,))
                conn.commit()
            finally:
                pass  # persistent connection

    def get_all_documents(self) -> List[Document]:
        conn = self._get_conn()
        try:
            rows = conn.execute("SELECT * FROM documents ORDER BY organized_at DESC").fetchall()
            return [Document(**dict(r)) for r in rows]
        finally:
            pass

    def get_documents_by_folder(self, folder: str) -> List[Document]:
        conn = self._get_conn()
        try:
            rows = conn.execute("SELECT * FROM documents WHERE target_folder=? ORDER BY organized_at DESC", (folder,)).fetchall()
            return [Document(**dict(r)) for r in rows]
        finally:
            pass

    def get_folders(self) -> List[str]:
        conn = self._get_conn()
        try:
            rows = conn.execute("SELECT DISTINCT target_folder FROM documents ORDER BY target_folder").fetchall()
            return [r['target_folder'] for r in rows]
        finally:
            pass

    def search_documents(self, keyword: str, search_content: bool = True) -> List[Document]:
        conn = self._get_conn()
        try:
            if search_content:
                rows = conn.execute(
                    "SELECT * FROM documents WHERE filename LIKE ? OR content_preview LIKE ? ORDER BY organized_at DESC",
                    (f"%{keyword}%", f"%{keyword}%")
                ).fetchall()
            else:
                rows = conn.execute(
                    "SELECT * FROM documents WHERE filename LIKE ? ORDER BY organized_at DESC",
                    (f"%{keyword}%",)
                ).fetchall()
            return [Document(**dict(r)) for r in rows]
        finally:
            pass

    def get_document_count(self) -> int:
        conn = self._get_conn()
        try:
            row = conn.execute("SELECT COUNT(*) as cnt FROM documents").fetchone()
            return row['cnt']
        finally:
            pass

    def get_folder_counts(self) -> Dict[str, int]:
        conn = self._get_conn()
        try:
            rows = conn.execute(
                "SELECT target_folder, COUNT(*) as cnt FROM documents GROUP BY target_folder ORDER BY cnt DESC"
            ).fetchall()
            return {r['target_folder']: r['cnt'] for r in rows}
        finally:
            pass

    # ---- 设置 ----
    def get_setting(self, key: str, default: str = "") -> str:
        conn = self._get_conn()
        try:
            row = conn.execute("SELECT value FROM settings WHERE key=?", (key,)).fetchone()
            return row['value'] if row else default
        finally:
            pass

    def set_setting(self, key: str, value: str):
        with self._lock:
            conn = self._get_conn()
            try:
                conn.execute("INSERT OR REPLACE INTO settings (key, value) VALUES (?,?)", (key, value))
                conn.commit()
            finally:
                pass  # persistent connection

    def get_documents_without_content(self) -> List[Document]:
        """获取未提取内容的文档"""
        conn = self._get_conn()
        try:
            rows = conn.execute(
                "SELECT * FROM documents WHERE content_preview IS NULL OR content_preview = '' ORDER BY organized_at DESC"
            ).fetchall()
            return [Document(**dict(r)) for r in rows]
        finally:
            pass

    def update_document_content(self, doc_id: int, content: str):
        """更新文档内容预览"""
        with self._lock:
            conn = self._get_conn()
            try:
                conn.execute("UPDATE documents SET content_preview = ? WHERE id = ?", (content, doc_id))
                conn.commit()
            finally:
                pass


# ==================== 内容提取器 ====================
class ContentExtractor:
    """从各种文档格式中提取文本内容"""

    @staticmethod
    def get_file_hash(filepath: str) -> str:
        """计算文件MD5哈希"""
        h = hashlib.md5()
        try:
            with open(filepath, 'rb') as f:
                for chunk in iter(lambda: f.read(8192), b''):
                    h.update(chunk)
            return h.hexdigest()
        except Exception:
            return ""

    @staticmethod
    def extract_text(filepath: str, max_chars: int = 2000) -> str:
        """
        提取文档文本内容（用于搜索索引）
        依次尝试：专用提取 → 通用备用提取，取最佳结果
        """
        ext = os.path.splitext(filepath)[1].lower()
        result = ""

        # 第一步：专用格式提取
        try:
            if ext in ('.txt', '.csv', '.md', '.log', '.json', '.xml', '.html', '.htm'):
                result = ContentExtractor._extract_text(filepath, max_chars)
            elif ext in ('.docx',):
                result = ContentExtractor._extract_docx(filepath, max_chars)
            elif ext in ('.doc', '.wps'):
                result = ContentExtractor._extract_doc(filepath, max_chars)
            elif ext in ('.xls', '.xlsx'):
                result = ContentExtractor._extract_excel(filepath, max_chars)
            elif ext in ('.et', '.ett'):
                result = ContentExtractor._extract_et(filepath, max_chars)
            elif ext == '.pdf':
                result = ContentExtractor._extract_pdf(filepath, max_chars)
            elif ext in ('.pptx',):
                result = ContentExtractor._extract_pptx(filepath, max_chars)
            elif ext in ('.ppt', '.dps'):
                result = ContentExtractor._extract_ppt(filepath, max_chars)
            elif ext == '.rtf':
                result = ContentExtractor._extract_rtf(filepath, max_chars)
        except Exception as e:
            logger.warning(f"提取内容失败 [{filepath}]: {e}")

        # 第二步：如果专用提取没拿到内容，尝试 OLE 提取（.doc/.wps/.et）
        if not result or not result.strip():
            if ext in ('.doc', '.wps', '.et', '.ett'):
                try:
                    import olefile
                    ole = olefile.OleFileIO(filepath)
                    try:
                        result = ContentExtractor._extract_ole_text(ole)
                    finally:
                        ole.close()
                except Exception:
                    pass

        # 第三步：如果仍然没有，尝试通用备用提取
        if not result or not result.strip():
            result = ContentExtractor._extract_binary_fallback(filepath, max_chars)

        return result

    @staticmethod
    def _extract_text(filepath: str, max_chars: int) -> str:
        try:
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read(max_chars)
        except:
            return ""

    @staticmethod
    def _extract_docx(filepath: str, max_chars: int) -> str:
        try:
            from docx import Document
            doc = Document(filepath)
            text = '\n'.join(p.text for p in doc.paragraphs if p.text.strip())
            return text[:max_chars]
        except ImportError:
            logger.warning("python-docx 未安装，无法提取 .docx 内容")
            return ""
        except Exception as e:
            logger.warning(f"提取docx失败: {e}")
            return ""

    @staticmethod
    def _extract_doc(filepath: str, max_chars: int) -> str:
        """提取 .doc/.wps 文件文本"""
        # 方法1：antiword（最准确）
        try:
            import subprocess
            result = subprocess.run(['antiword', filepath], capture_output=True, text=True, timeout=10)
            if result.returncode == 0 and result.stdout.strip():
                return result.stdout[:max_chars]
        except Exception:
            pass
        # 方法2：catdoc
        try:
            import subprocess
            result = subprocess.run(['catdoc', '-w', filepath], capture_output=True, text=True, timeout=10)
            if result.returncode == 0 and result.stdout.strip():
                return result.stdout[:max_chars]
        except Exception:
            pass
        # 方法3：olefile 读取 OLE 复合文档中的文本
        try:
            import olefile
            ole = olefile.OleFileIO(filepath)
            try:
                text = ContentExtractor._extract_ole_text(ole)
                if text and text.strip():
                    return text[:max_chars]
            finally:
                ole.close()
        except ImportError:
            pass
        except Exception:
            pass
        return ""

    @staticmethod
    def _extract_ole_text(ole) -> str:
        """从 OLE 复合文档中提取文本（支持 .doc/.wps）"""
        texts = []

        # .doc 文本主要在 WordDocument 流中
        # 文本以 UTF-16LE 编码存储（现代 .doc）或 CP1252（旧 .doc）
        if ole.exists('WordDocument'):
            try:
                stream = ole.openstream('WordDocument')
                data = stream.read()

                # .doc 的 FIB（File Information Block）结构：
                # 偏移 0x0000: wIdent (0xA5EC = .doc)
                # 偏移 0x000C: fcClx (指向 CLX 的偏移量，相对于 WordDocument stream)
                # 偏移 0x0018: ccpText (文档正文文本的字符数)
                # 偏移 0x004A: fcMin (文本起始偏移量，相对于 WordDocument stream)
                # 偏移 0x004E: fcMac (文本结束偏移量)

                if len(data) > 0x4F:
                    import struct
                    wIdent = struct.unpack_from('<H', data, 0)[0]
                    if wIdent == 0xA5EC:  # 确认是 .doc 格式
                        # 尝试读取文本范围
                        fcMin = struct.unpack_from('<I', data, 0x004A)[0]
                        fcMac = struct.unpack_from('<I', data, 0x004E)[0]
                        ccpText = struct.unpack_from('<I', data, 0x0018)[0]

                        if 0 < fcMin < fcMac <= len(data) and ccpText > 0:
                            # 文本在 WordDocument 流的 fcMin 到 fcMac 之间
                            # 以 UTF-16LE 编码
                            raw_text = data[fcMin:fcMac]
                            try:
                                text = raw_text.decode('utf-16-le', errors='ignore')
                                # 去除 .doc 内部控制字符（如 \r\x07 = 表格单元格结束, \x13 = 域开始等）
                                text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x0b]', '', text)
                                # 去除 .doc 特殊标记
                                text = text.replace('\r\x07', '\n')  # 单元格结束
                                text = text.replace('\r\x13', '')     # 域开始
                                text = text.replace('\x14', '')       # 域分隔
                                text = text.replace('\x15', '')       # 域结束
                                text = text.replace('\x01', '')       # 图片/对象标记
                                text = text.replace('\x08', '')       # 不间断空格标记
                                if text.strip():
                                    texts.append(text)
                            except Exception:
                                pass

                # 如果 FIB 解析失败，尝试直接从流中扫描文本
                if not texts:
                    # 尝试 UTF-16LE 解码整个流
                    try:
                        text = data.decode('utf-16-le', errors='ignore')
                        clean = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', '', text)
                        segments = re.findall(r'[\u4e00-\u9fff\u3000-\u303f\uff00-\uffefa-zA-Z0-9\s,.;:!?，。；：！？、（）]{4,}', clean)
                        if segments:
                            texts.append('\n'.join(segments))
                    except Exception:
                        pass

                # 尝试 GBK/GB2312 解码（某些中文 .doc）
                if not texts:
                    try:
                        text = data.decode('gbk', errors='ignore')
                        clean = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', '', text)
                        segments = re.findall(r'[\u4e00-\u9fff\u3000-\u303f\uff00-\uffefa-zA-Z0-9\s,.;:!?，。；：！？、（）]{4,}', clean)
                        if segments:
                            texts.append('\n'.join(segments))
                    except Exception:
                        pass
            except Exception:
                pass

        # 扫描所有流中的文本（.wps 可能有不同的流名称）
        if not texts:
            for stream_name in ole.listdir():
                stream_path = '/'.join(stream_name)
                # 跳过已处理的和非文本流
                if any(k in stream_path.lower() for k in (
                    'compobj', 'summary', 'properties', '\x01ole',
                    'preview', 'thumbnail', 'encryption', 'revision'
                )):
                    continue
                try:
                    data = ole.openstream(stream_name).read()
                    if len(data) < 10:
                        continue
                    # 尝试多种编码
                    for encoding in ('utf-16-le', 'gbk', 'gb2312', 'gb18030', 'utf-8'):
                        try:
                            text = data.decode(encoding, errors='ignore')
                            clean = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', '', text)
                            segments = re.findall(r'[\u4e00-\u9fff\u3000-\u303f\uff00-\uffefa-zA-Z0-9\s,.;:!?，。；：！？、（）]{4,}', clean)
                            if segments and len('\n'.join(segments)) > 20:
                                texts.append('\n'.join(segments))
                                break
                        except Exception:
                            continue
                except Exception:
                    continue

        return '\n'.join(texts) if texts else ""

    @staticmethod
    def _extract_excel(filepath: str, max_chars: int) -> str:
        try:
            import openpyxl
            wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
            texts = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(max_row=100, values_only=True):
                    row_text = ' '.join(str(c) for c in row if c is not None)
                    if row_text.strip():
                        texts.append(row_text)
            wb.close()
            return '\n'.join(texts)[:max_chars]
        except ImportError:
            logger.warning("openpyxl 未安装，无法提取 Excel 内容")
            return ""
        except Exception as e:
            logger.warning(f"提取Excel失败: {e}")
            return ""

    @staticmethod
    def _extract_pdf(filepath: str, max_chars: int) -> str:
        # 优先使用 pdfplumber（中文提取效果更好）
        try:
            import pdfplumber
            texts = []
            with pdfplumber.open(filepath) as pdf:
                for page in pdf.pages[:20]:
                    t = page.extract_text()
                    if t:
                        texts.append(t)
            result = '\n'.join(texts)
            if result.strip():
                # 检测乱码：如果中文字符占比过低，可能是编码问题
                if ContentExtractor._is_garbled(result):
                    pass  # fallback 到下一个方法
                else:
                    return result[:max_chars]
        except ImportError:
            pass
        except Exception:
            pass
        # 备用：PyPDF2
        try:
            import PyPDF2
            text = ""
            with open(filepath, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages[:20]:
                    t = page.extract_text()
                    if t:
                        text += t + "\n"
            if text.strip():
                if ContentExtractor._is_garbled(text):
                    pass  # fallback
                else:
                    return text[:max_chars]
        except ImportError:
            pass
        except Exception:
            pass
        # 无依赖：直接从 PDF 二进制中提取文本流
        return ContentExtractor._extract_pdf_raw(filepath, max_chars)

    @staticmethod
    def _is_garbled(text: str) -> bool:
        """检测文本是否为乱码（中文字符占比过低）"""
        if not text:
            return True
        total = len(text)
        if total < 10:
            return False
        # 统计中文字符数量
        chinese = sum(1 for c in text if '\u4e00' <= c <= '\u9fff')
        # 统计可读 ASCII 字符
        ascii_readable = sum(1 for c in text if c.isalnum() or c in '，。；：！？、（）""''【】《》 \n')
        readable_ratio = (chinese + ascii_readable) / total
        return readable_ratio < 0.3

    @staticmethod
    def _extract_pdf_raw(filepath: str, max_chars: int) -> str:
        """无依赖 PDF 文本提取：解析 PDF 中的文本流"""
        try:
            with open(filepath, 'rb') as f:
                data = f.read()

            texts = []

            # 方法1：提取 BT...ET 块中的文本
            bt_pattern = re.compile(rb'BT(.*?)ET', re.DOTALL)
            for bt_match in bt_pattern.finditer(data):
                block = bt_match.group(1)

                # 处理 (text) Tj 格式 — PDFDocEncoding
                for m in re.finditer(rb'\(([^)]*)\)\s*Tj', block):
                    raw = m.group(1)
                    text = ContentExtractor._decode_pdf_string(raw)
                    if text and text.strip():
                        texts.append(text)

                # 处理 <hex> Tj 格式 — 十六进制编码
                for m in re.finditer(rb'<([0-9A-Fa-f]+)>\s*Tj', block):
                    hex_str = m.group(1).decode('ascii', errors='ignore')
                    text = ContentExtractor._decode_pdf_hex(hex_str)
                    if text and text.strip():
                        texts.append(text)

                # 处理原始字节 <bytes> Tj 格式（部分PDF直接嵌入UTF-16BE字节）
                for m in re.finditer(rb'<([\x80-\xff][\x00-\xff]*)>\s*Tj', block):
                    raw = m.group(1)
                    text = ContentExtractor._decode_pdf_bytes(raw)
                    if text and text.strip():
                        texts.append(text)

                # 处理 TJ 数组格式 [...]
                for m in re.finditer(rb'\[(.*?)\]\s*TJ', block, re.DOTALL):
                    array_content = m.group(1)
                    line = ""
                    for part in re.findall(rb'\(([^)]*)\)', array_content):
                        t = ContentExtractor._decode_pdf_string(part)
                        if t:
                            line += t
                    for hex_str in re.findall(rb'<([0-9A-Fa-f]+)>', array_content):
                        t = ContentExtractor._decode_pdf_hex(hex_str)
                        if t:
                            line += t
                    for m2 in re.finditer(rb'<([\x80-\xff][\x00-\xff]*)>', array_content):
                        t = ContentExtractor._decode_pdf_bytes(m2.group(1))
                        if t:
                            line += t
                    if line.strip():
                        texts.append(line)

            if texts:
                return '\n'.join(texts)[:max_chars]

            # 方法2：在整个文件中扫描十六进制和原始字节编码的中文
            hex_texts = []
            for m in re.finditer(rb'<([0-9A-Fa-f]{8,})>', data):
                text = ContentExtractor._decode_pdf_hex(m.group(1).decode('ascii', errors='ignore'))
                if text:
                    clean = re.findall(r'[\u4e00-\u9fff\u3000-\u303f\uff00-\uffefa-zA-Z0-9]{2,}', text)
                    if clean:
                        hex_texts.extend(clean)
            for m in re.finditer(rb'<([\x80-\xff][\x00-\xff]{4,})>', data):
                text = ContentExtractor._decode_pdf_bytes(m.group(1))
                if text:
                    clean = re.findall(r'[\u4e00-\u9fff\u3000-\u303f\uff00-\uffefa-zA-Z0-9]{2,}', text)
                    if clean:
                        hex_texts.extend(clean)
            if hex_texts:
                return '\n'.join(hex_texts)[:max_chars]

            # 方法3：直接在整个文件中搜索文本字符串
            for encoding in ('utf-16-be', 'latin-1'):
                try:
                    text = data.decode(encoding, errors='ignore')
                    segments = re.findall(
                        r'[\u4e00-\u9fff\u3000-\u303f\uff00-\uffef]{2,}|'
                        r'[a-zA-Z0-9\s,.;:!?]{4,}',
                        text
                    )
                    if segments:
                        return '\n'.join(segments)[:max_chars]
                except Exception:
                    continue

        except Exception:
            pass
        return ""

    @staticmethod
    def _decode_pdf_string(raw: bytes) -> str:
        """解码 PDF () 格式字符串"""
        try:
            try:
                text = raw.decode('utf-16-be')
            except Exception:
                text = raw.decode('latin-1')
            text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', '', text)
            return text
        except Exception:
            return ""

    @staticmethod
    def _decode_pdf_hex(hex_str: str) -> str:
        """解码 PDF <> 十六进制格式字符串"""
        try:
            raw_bytes = bytes.fromhex(hex_str)
            return ContentExtractor._decode_pdf_bytes(raw_bytes)
        except Exception:
            return ""

    @staticmethod
    def _decode_pdf_bytes(raw_bytes: bytes) -> str:
        """解码 PDF 原始字节（自动检测编码）"""
        try:
            if len(raw_bytes) < 2:
                return ""
            # 检测 BOM
            if raw_bytes[:2] == b'\xfe\xff':
                text = raw_bytes[2:].decode('utf-16-be', errors='ignore')
            elif raw_bytes[:2] == b'\xff\xfe':
                text = raw_bytes[2:].decode('utf-16-le', errors='ignore')
            else:
                # 尝试 UTF-16BE（中文PDF最常用）
                try:
                    text = raw_bytes.decode('utf-16-be', errors='ignore')
                    # 验证解码质量：如果中文字符占比太低，可能不是UTF-16BE
                    chinese = sum(1 for c in text if '\u4e00' <= c <= '\u9fff')
                    if chinese > 0:
                        return re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', '', text)
                except Exception:
                    pass
                text = raw_bytes.decode('latin-1', errors='ignore')
            return re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', '', text)
        except Exception:
            return ""

    @staticmethod
    def _extract_pptx(filepath: str, max_chars: int) -> str:
        try:
            from pptx import Presentation
            prs = Presentation(filepath)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            if paragraph.text.strip():
                                texts.append(paragraph.text)
            return '\n'.join(texts)[:max_chars]
        except ImportError:
            logger.warning("python-pptx 未安装，无法提取 .pptx 内容")
            return ""
        except Exception as e:
            logger.warning(f"提取pptx失败: {e}")
            return ""

    @staticmethod
    def _extract_ppt(filepath: str, max_chars: int) -> str:
        """旧版PPT格式，尝试基本提取"""
        try:
            import subprocess
            result = subprocess.run(['catppt', filepath], capture_output=True, text=True, timeout=10)
            if result.returncode == 0:
                return result.stdout[:max_chars]
        except Exception:
            pass
        return ""

    @staticmethod
    def _extract_et(filepath: str, max_chars: int) -> str:
        """提取 WPS 表格 (.et/.ett) 内容"""
        # .et 本质是 ZIP 格式（类似 .xlsx），尝试用 openpyxl 读取
        try:
            import openpyxl
            wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
            texts = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(max_row=100, values_only=True):
                    row_text = ' '.join(str(c) for c in row if c is not None)
                    if row_text.strip():
                        texts.append(row_text)
            wb.close()
            return '\n'.join(texts)[:max_chars]
        except Exception:
            pass
        # 尝试作为 OLE 复合文档读取
        try:
            import olefile
            ole = olefile.OleFileIO(filepath)
            texts = []
            for stream_name in ole.listdir():
                stream_path = '/'.join(stream_name)
                try:
                    data = ole.openstream(stream_name).read()
                    text = data.decode('utf-8', errors='ignore')
                    # 提取可读文本
                    clean = re.sub(r'[^\u4e00-\u9fff\u3000-\u303f\uff00-\uffefa-zA-Z0-9\s,.;:!?，。；：！？、（）\n]', '', text)
                    if len(clean.strip()) > 5:
                        texts.append(clean.strip())
                except Exception:
                    continue
            ole.close()
            return '\n'.join(texts)[:max_chars]
        except Exception:
            pass
        return ""

    @staticmethod
    def _extract_rtf(filepath: str, max_chars: int) -> str:
        """提取 RTF 文件中的纯文本"""
        try:
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                data = f.read(max_chars * 4)
            # 去除 RTF 控制词，保留纯文本
            text = re.sub(r'\\[a-z]+\d* ? ?', '', data)
            text = re.sub(r'[{}]', '', text)
            text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x9f]', '', text)
            # 提取连续可读文本
            segments = re.findall(r'[\u4e00-\u9fff\u3000-\u303f\uff00-\uffefa-zA-Z0-9\s,.;:!?，。；：！？]{4,}', text)
            return '\n'.join(segments)[:max_chars]
        except Exception:
            return ""

    @staticmethod
    def _extract_binary_fallback(filepath: str, max_chars: int) -> str:
        """通用备用提取：从二进制文件中提取可读文本"""
        try:
            # 先尝试作为 ZIP 文件提取（docx/xlsx/pptx/et 本质都是 ZIP）
            if zipfile.is_zipfile(filepath):
                text = ContentExtractor._extract_from_zip(filepath, max_chars)
                if text and text.strip():
                    return text[:max_chars]

            with open(filepath, 'rb') as f:
                data = f.read(max_chars * 4)
            # 使用 chardet 检测编码
            try:
                import chardet
                detected = chardet.detect(data)
                encoding = detected.get('encoding', 'utf-8') if detected else 'utf-8'
                if encoding and encoding.lower() not in ('ascii', 'utf-8', 'utf-16-le'):
                    try:
                        text = data.decode(encoding, errors='ignore')
                        clean = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', '', text)
                        segments = re.findall(r'[\u4e00-\u9fff\u3000-\u303f\uff00-\uffefa-zA-Z0-9\s,.;:!?，。；：！？、（）\n]{4,}', clean)
                        result = '\n'.join(segments)
                        if len(result) > 20:
                            return result[:max_chars]
                    except Exception:
                        pass
            except ImportError:
                pass
            # 尝试多种编码
            for encoding in ('utf-8', 'gbk', 'gb2312', 'gb18030', 'big5', 'utf-16-le'):
                try:
                    text = data.decode(encoding, errors='ignore')
                    clean = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', '', text)
                    segments = re.findall(r'[\u4e00-\u9fff\u3000-\u303f\uff00-\uffefa-zA-Z0-9\s,.;:!?，。；：！？、（）""''【】《》\-_]{4,}', clean)
                    result = '\n'.join(segments)
                    if len(result) > 20:
                        return result[:max_chars]
                except Exception:
                    continue
        except Exception:
            pass
        return ""

    @staticmethod
    def _extract_from_zip(filepath: str, max_chars: int) -> str:
        """从 ZIP 格式文件中提取文本（docx/xlsx/pptx/et 等都是 ZIP）"""
        import zipfile
        import html
        texts = []
        try:
            with zipfile.ZipFile(filepath, 'r') as zf:
                for name in zf.namelist():
                    # docx: word/document.xml, xlsx: xl/sharedStrings.xml, xl/worksheets/sheet*.xml
                    # pptx: ppt/slides/slide*.xml
                    name_lower = name.lower()
                    should_extract = any(k in name_lower for k in (
                        'document.xml', 'sharedstrings', 'slide', 'notes',
                        'worksheet', 'sheet',
                        '.txt', '.md', '.csv', '.json', '.html'
                    ))
                    # 跳过样式/主题等无意义文件
                    if any(k in name_lower for k in ('style', 'theme', 'font', 'setting', 'layout')):
                        continue
                    if not should_extract:
                        continue

                    try:
                        data = zf.read(name)
                        text = data.decode('utf-8', errors='ignore')

                        if name_lower.endswith('.xml') or name_lower.endswith('.rels'):
                            # 解码 HTML 实体（如 &amp; -> &）
                            text = html.unescape(text)
                            # 去除 XML 标签
                            clean = re.sub(r'<[^>]+>', '\n', text)
                            # 去除多余空白
                            clean = re.sub(r'\n\s*\n', '\n', clean)
                            clean = clean.strip()
                            if len(clean) > 2:
                                texts.append(clean)
                        else:
                            if len(text.strip()) > 5:
                                texts.append(text)
                    except Exception:
                        continue
        except Exception:
            pass
        return '\n'.join(texts)[:max_chars]


# ==================== 依赖诊断 ====================
def check_dependencies() -> dict:
    """检查内容提取依赖安装状态，返回 {模块名: 是否可用}"""
    deps = {
        'python-docx (Word .docx)': False,
        'openpyxl (Excel .xlsx)': False,
        'PyPDF2 (PDF)': False,
        'pdfplumber (PDF增强)': False,
        'python-pptx (PPT .pptx)': False,
        'olefile (旧版Office)': False,
        'chardet (编码检测)': False,
    }
    try:
        import docx; deps['python-docx (Word .docx)'] = True
    except ImportError: pass
    try:
        import openpyxl; deps['openpyxl (Excel .xlsx)'] = True
    except ImportError: pass
    try:
        import PyPDF2; deps['PyPDF2 (PDF)'] = True
    except ImportError: pass
    try:
        import pdfplumber; deps['pdfplumber (PDF增强)'] = True
    except ImportError: pass
    try:
        import pptx; deps['python-pptx (PPT .pptx)'] = True
    except ImportError: pass
    try:
        import olefile; deps['olefile (旧版Office)'] = True
    except ImportError: pass
    try:
        import chardet; deps['chardet (编码检测)'] = True
    except ImportError: pass
    return deps


def test_extract(filepath: str) -> str:
    """测试单个文件的内容提取，返回提取到的文本"""
    if not os.path.isfile(filepath):
        return f"文件不存在: {filepath}"
    text = ContentExtractor.extract_text(filepath)
    if text and text.strip():
        return f"提取成功 ({len(text)} 字符):\n{text[:500]}"
    else:
        return "提取失败：未获取到文本内容。请检查文件是否损坏或为空文件。"


# ==================== 全盘搜索 ====================
class GlobalSearcher:
    """全盘文件搜索器（支持实时回调、用户指定路径）"""

    # 搜索时跳过的目录（仅顶级系统目录）
    SKIP_DIRS = {
        '__pycache__', 'node_modules', '.git', 'venv', 'env', '.venv',
        '$RECYCLE.BIN', 'System Volume Information', 'lost+found',
        'proc', 'sys', 'dev', 'run', 'snap', 'boot',
        'Windows', 'Program Files', 'Program Files (x86)', 'ProgramData',
        'AppData', 'Application Data',
    }

    def __init__(self):
        self._stop = threading.Event()

    def search(self, keyword: str, search_paths: List[str] = None,
               search_content: bool = False,
               on_found=None, on_progress=None, on_done=None):
        """
        在指定路径中搜索文件（实时回调）
        keyword: 搜索关键词
        search_paths: 搜索路径列表
        search_content: 是否搜索文件内容
        on_found: 找到文件时回调 on_found(doc) — 实时显示
        on_progress: 进度回调 on_progress(searched_count, found_count, current_path)
        on_done: 搜索完成回调 on_done(total_results)
        """
        self._stop.clear()
        keyword_lower = keyword.lower()
        total_searched = 0
        total_found = 0

        if not search_paths:
            search_paths = self._get_all_available_paths()

        for search_path in search_paths:
            if self._stop.is_set():
                break
            if not os.path.isdir(search_path):
                continue

            try:
                for root, dirs, files in os.walk(search_path, topdown=True):
                    if self._stop.is_set():
                        break

                    # 跳过系统/隐藏目录
                    dirs[:] = [d for d in dirs
                               if not d.startswith('.') and d not in self.SKIP_DIRS]

                    for filename in files:
                        if self._stop.is_set():
                            break

                        filepath = os.path.join(root, filename)
                        ext = os.path.splitext(filename)[1].lower()

                        # 按文件名匹配
                        name_match = keyword_lower in filename.lower()

                        # 按内容匹配（对所有可提取文本的格式）
                        content_text = ""
                        content_match = False
                        if search_content and not name_match:
                            # 对所有支持格式和纯文本文件做内容搜索
                            if ext in TEXT_EXTRACTABLE_EXTENSIONS or ext in (
                                '.txt', '.md', '.log', '.csv', '.json', '.xml',
                                '.html', '.htm', '.rtf', '.ini', '.cfg', '.conf',
                                '.py', '.java', '.c', '.cpp', '.h', '.js', '.sh',
                                '.bat', '.ps1', '.yml', '.yaml', '.toml', '.properties',
                            ):
                                try:
                                    content_text = ContentExtractor.extract_text(filepath)
                                    if content_text and keyword_lower in content_text.lower():
                                        content_match = True
                                except Exception:
                                    # 加密文件、损坏文件等：跳过内容搜索，不中断整个搜索
                                    content_text = ""

                        if name_match or content_match:
                            try:
                                file_size = os.path.getsize(filepath)
                            except Exception:
                                file_size = 0

                            doc = Document(
                                filename=filename,
                                filepath=filepath,
                                target_folder=root,
                                rule_name="全盘搜索",
                                file_size=file_size,
                                file_type=EXTENSION_TYPE_MAP.get(ext, "其他"),
                                content_preview=content_text[:2000] if content_text else "",
                                organized_at=datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                                original_path=filepath
                            )
                            total_found += 1
                            if on_found:
                                try:
                                    on_found(doc)
                                except Exception:
                                    pass

                        total_searched += 1
                        if on_progress and total_searched % 20 == 0:
                            try:
                                on_progress(total_searched, total_found, root)
                            except Exception:
                                pass
            except PermissionError:
                continue
            except Exception:
                continue

        if on_done:
            try:
                on_done(total_found)
            except Exception:
                pass

    def stop(self):
        self._stop.set()

    @staticmethod
    def _get_all_available_paths() -> List[str]:
        """获取所有可搜索的路径"""
        paths = []

        if IS_WINDOWS:
            # Windows: 所有可用盘符
            for drive in "CDEFGHIJKLMNOPQRSTUVWXYZ":
                drive_path = f"{drive}:\\"
                if os.path.isdir(drive_path):
                    paths.append(drive_path)
        elif IS_LINUX:
            # Linux/麒麟: 常见挂载点
            for p in ("/home", "/data", "/media", "/mnt", "/tmp", "/root"):
                if os.path.isdir(p):
                    paths.append(p)
            # 读取实际磁盘挂载点（排除虚拟文件系统）
            virtual_fs = {'proc', 'sys', 'dev', 'run', 'tmpfs', 'cgroup', 'debugfs',
                          'configfs', 'fusectl', 'mqueue', 'hugetlbfs', 'pstore',
                          'binfmt_misc', 'nfsd', 'rpc_pipefs'}
            try:
                with open('/proc/mounts', 'r') as f:
                    for line in f:
                        parts = line.split()
                        if len(parts) >= 3:
                            mount = parts[1]
                            fstype = parts[2]
                            if (mount not in paths and mount.startswith('/')
                                    and os.path.isdir(mount)
                                    and fstype not in virtual_fs
                                    and not mount.startswith('/proc')
                                    and not mount.startswith('/sys')
                                    and not mount.startswith('/dev')):
                                paths.append(mount)
            except Exception:
                pass
        elif IS_MACOS:
            paths.append("/")
            paths.append(os.path.expanduser("~"))

        return paths if paths else [os.path.expanduser("~")]


# ==================== 规则引擎 ====================
class RuleEngine:
    """规则匹配引擎"""

    @staticmethod
    def match(filename: str, rule: Rule) -> bool:
        """
        判断文件名是否匹配规则
        contain_keywords: 所有关键词都必须包含（AND关系）
        exclude_keywords: 任一关键词匹配则排除（OR关系）
        """
        if not rule.enabled:
            return False

        name = os.path.splitext(filename)[0]  # 不含扩展名

        # 检查包含关键词
        if rule.contain_keywords:
            keywords = [k.strip() for k in rule.contain_keywords.split(',') if k.strip()]
            if keywords:
                for kw in keywords:
                    if kw not in name:
                        return False

        # 检查排除关键词
        if rule.exclude_keywords:
            keywords = [k.strip() for k in rule.exclude_keywords.split(',') if k.strip()]
            for kw in keywords:
                if kw in name:
                    return False

        return True

    @staticmethod
    def find_matching_rule(filename: str, rules: List[Rule]) -> Optional[Rule]:
        """按优先级找到第一个匹配的规则"""
        for rule in rules:
            if RuleEngine.match(filename, rule):
                return rule
        return None


# ==================== 文件整理器 ====================
class FileOrganizer:
    """文件整理核心逻辑"""

    def __init__(self, db: Database):
        self.db = db
        self._stop_event = threading.Event()

    def get_desktop_path(self) -> str:
        """获取桌面路径（跨平台）"""
        desktop = self.db.get_setting("desktop_path", "")
        if desktop and os.path.isdir(desktop):
            return desktop

        home = os.path.expanduser("~")

        if IS_WINDOWS:
            # Windows: 尝试 Desktop 和 桌面
            for name in ("Desktop", "桌面"):
                path = os.path.join(home, name)
                if os.path.isdir(path):
                    return path
            # 尝试通过注册表获取（中文Windows桌面可能在其他位置）
            try:
                import winreg
                key = winreg.OpenKey(winreg.HKEY_CURRENT_USER,
                                     r"Software\Microsoft\Windows\CurrentVersion\Explorer\User Shell Folders")
                desktop_val = winreg.QueryValueEx(key, "Desktop")[0]
                winreg.CloseKey(key)
                # 处理 %USERPROFILE% 等环境变量
                desktop_val = os.path.expandvars(desktop_val)
                if os.path.isdir(desktop_val):
                    return desktop_val
            except Exception:
                pass
        elif IS_LINUX:
            # Linux/麒麟: 优先使用 XDG 标准路径
            xdg_desktop = os.environ.get("XDG_DESKTOP_DIR", "")
            if xdg_desktop and os.path.isdir(xdg_desktop):
                return xdg_desktop
            # 常见桌面路径
            for name in ("Desktop", "桌面"):
                path = os.path.join(home, name)
                if os.path.isdir(path):
                    return path
            # 麒麟系统可能使用中文桌面
            for name in ("桌面", "desktop"):
                path = os.path.join(home, name)
                if os.path.isdir(path):
                    return path
        elif IS_MACOS:
            path = os.path.join(home, "Desktop")
            if os.path.isdir(path):
                return path

        return ""

    def scan_desktop(self) -> List[str]:
        """扫描桌面上的支持格式文件"""
        desktop = self.get_desktop_path()
        if not desktop:
            return []
        files = []
        for f in os.listdir(desktop):
            ext = os.path.splitext(f)[1].lower()
            if ext in SUPPORTED_EXTENSIONS:
                files.append(os.path.join(desktop, f))
        return files

    def organize_files(self, progress_callback=None) -> Tuple[int, int, List[str]]:
        """
        一键整理桌面文件
        返回: (成功数, 跳过数, 错误信息列表)
        """
        rules = self.db.get_enabled_rules()
        if not rules:
            return 0, 0, ["没有可用的分类规则，请先添加规则"]

        files = self.scan_desktop()
        if not files:
            return 0, 0, ["桌面上没有找到支持的文档文件"]

        success_count = 0
        skip_count = 0
        errors = []
        total = len(files)

        for i, filepath in enumerate(files):
            if self._stop_event.is_set():
                errors.append("操作已被用户取消")
                break

            filename = os.path.basename(filepath)
            if progress_callback:
                progress_callback(i + 1, total, f"正在处理: {filename}")

            try:
                rule = RuleEngine.find_matching_rule(filename, rules)
                if not rule:
                    skip_count += 1
                    logger.info(f"跳过（无匹配规则）: {filename}")
                    continue

                # 检查目标文件夹是否存在，不存在则创建
                target_dir = rule.target_folder
                if not os.path.exists(target_dir):
                    os.makedirs(target_dir, exist_ok=True)
                    logger.info(f"创建文件夹: {target_dir}")

                # 检查目标是否已存在同名文件
                target_path = os.path.join(target_dir, filename)
                if os.path.exists(target_path):
                    # 添加时间戳避免覆盖
                    stem = os.path.splitext(filename)[0]
                    ext = os.path.splitext(filename)[1]
                    timestamp = datetime.now().strftime("_%H%M%S")
                    filename = f"{stem}{timestamp}{ext}"
                    target_path = os.path.join(target_dir, filename)

                # 移动文件
                shutil.move(filepath, target_path)

                # 提取内容用于搜索索引
                content = ContentExtractor.extract_text(target_path)
                file_hash = ContentExtractor.get_file_hash(target_path)
                file_size = os.path.getsize(target_path)
                file_type = EXTENSION_TYPE_MAP.get(os.path.splitext(filename)[1].lower(), "未知")

                # 记录到数据库
                doc = Document(
                    filename=filename,
                    filepath=target_path,
                    target_folder=target_dir,
                    rule_name=rule.name,
                    file_size=file_size,
                    file_type=file_type,
                    content_preview=content,
                    file_hash=file_hash,
                    organized_at=datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                    original_path=filepath
                )
                self.db.add_document(doc)
                success_count += 1
                logger.info(f"已整理: {filename} → {target_dir}")

            except Exception as e:
                errors.append(f"处理 {filename} 失败: {str(e)}")
                logger.error(f"处理失败 [{filename}]: {e}")

        return success_count, skip_count, errors

    def stop(self):
        """停止整理操作"""
        self._stop_event.set()

    def reset_stop(self):
        """重置停止标志"""
        self._stop_event.clear()


# ==================== 实时监控器 ====================
class DesktopWatcher:
    """桌面文件实时监控"""

    def __init__(self, db: Database, callback=None):
        self.db = db
        self.callback = callback
        self._watching = False
        self._thread = None
        self._organizer = FileOrganizer(db)
        self._processed_files = set()
        self._scan_interval = 3  # 秒

    def start(self):
        """开始监控"""
        if self._watching:
            return
        self._watching = True
        # 初始化已处理文件列表
        for f in self._organizer.scan_desktop():
            self._processed_files.add(os.path.basename(f))
        self._thread = threading.Thread(target=self._watch_loop, daemon=True)
        self._thread.start()
        logger.info("桌面监控已启动")

    def stop(self):
        """停止监控"""
        self._watching = False
        if self._thread:
            self._thread.join(timeout=5)
        logger.info("桌面监控已停止")

    def is_watching(self) -> bool:
        return self._watching

    def _watch_loop(self):
        while self._watching:
            try:
                files = self._organizer.scan_desktop()
                new_files = [f for f in files if os.path.basename(f) not in self._processed_files]
                if new_files:
                    logger.info(f"检测到 {len(new_files)} 个新文件")
                    rules = self.db.get_enabled_rules()
                    for filepath in new_files:
                        filename = os.path.basename(filepath)
                        rule = RuleEngine.find_matching_rule(filename, rules)
                        if rule:
                            try:
                                target_dir = rule.target_folder
                                if not os.path.exists(target_dir):
                                    os.makedirs(target_dir, exist_ok=True)
                                target_path = os.path.join(target_dir, filename)
                                if os.path.exists(target_path):
                                    stem = os.path.splitext(filename)[0]
                                    ext = os.path.splitext(filename)[1]
                                    timestamp = datetime.now().strftime("_%H%M%S")
                                    new_name = f"{stem}{timestamp}{ext}"
                                    target_path = os.path.join(target_dir, new_name)

                                shutil.move(filepath, target_path)
                                content = ContentExtractor.extract_text(target_path)
                                file_hash = ContentExtractor.get_file_hash(target_path)
                                file_size = os.path.getsize(target_path)
                                file_type = EXTENSION_TYPE_MAP.get(os.path.splitext(filename)[1].lower(), "未知")

                                doc = Document(
                                    filename=filename,
                                    filepath=target_path,
                                    target_folder=target_dir,
                                    rule_name=rule.name,
                                    file_size=file_size,
                                    file_type=file_type,
                                    content_preview=content,
                                    file_hash=file_hash,
                                    organized_at=datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                                    original_path=filepath
                                )
                                self.db.add_document(doc)
                                logger.info(f"自动整理: {filename} → {target_dir}")
                                if self.callback:
                                    self.callback(filename, target_dir, True, "")
                            except Exception as e:
                                logger.error(f"自动整理失败 [{filename}]: {e}")
                                if self.callback:
                                    self.callback(filename, "", False, str(e))
                        self._processed_files.add(filename)
            except Exception as e:
                logger.error(f"监控出错: {e}")
            time.sleep(self._scan_interval)


# ==================== 工具函数 ====================
def format_size(size_bytes: int) -> str:
    """格式化文件大小"""
    if size_bytes < 1024:
        return f"{size_bytes} B"
    elif size_bytes < 1024 * 1024:
        return f"{size_bytes / 1024:.1f} KB"
    elif size_bytes < 1024 * 1024 * 1024:
        return f"{size_bytes / (1024 * 1024):.1f} MB"
    else:
        return f"{size_bytes / (1024 * 1024 * 1024):.1f} GB"


def get_file_icon(file_type: str) -> str:
    """获取文件类型图标（emoji）"""
    icons = {
        'Word': '📝', 'WPS': '📝',
        'Excel': '📊', 'CSV': '📊',
        'PDF': '📄',
        'PPT': '📽️',
    }
    return icons.get(file_type, '📁')


def open_file(filepath: str):
    """跨平台打开文件"""
    if not os.path.exists(filepath):
        return False
    try:
        if IS_WINDOWS:
            os.startfile(filepath)
        elif IS_MACOS:
            subprocess.Popen(['open', filepath])
        else:
            # Linux/麒麟: 使用 xdg-open
            subprocess.Popen(['xdg-open', filepath])
        return True
    except Exception as e:
        logger.error(f"打开文件失败 [{filepath}]: {e}")
        return False


def locate_file(filepath: str):
    """跨平台在文件管理器中定位文件"""
    if not os.path.exists(filepath):
        return False
    try:
        if IS_WINDOWS:
            subprocess.Popen(f'explorer /select,"{filepath}"')
        elif IS_MACOS:
            subprocess.Popen(['open', '-R', filepath])
        else:
            # Linux/麒麟: 尝试多种文件管理器
            folder = os.path.dirname(filepath)
            # 尝试 dbus-send 打开文件管理器并选中文件
            try:
                subprocess.Popen([
                    'dbus-send', '--session', '--dest=org.freedesktop.FileManager1',
                    '--type=method_call', '/org/freedesktop/FileManager1',
                    'org.freedesktop.FileManager1.ShowItems',
                    f'array:string:file://{filepath}', 'string:'
                ])
            except Exception:
                # 回退: 直接打开所在文件夹
                subprocess.Popen(['xdg-open', folder])
        return True
    except Exception as e:
        logger.error(f"定位文件失败 [{filepath}]: {e}")
        return False


def select_folder_dialog(title: str = "选择文件夹", parent=None, initialdir=None) -> str:
    """
    弹出文件夹选择对话框（自定义实现，兼容麒麟系统）
    用户可以手动输入路径或浏览选择，点击「确定」返回路径
    """
    try:
        import tkinter as tk
        from tkinter import filedialog

        # 先尝试系统原生对话框
        try:
            if parent:
                folder = filedialog.askdirectory(title=title, parent=parent, initialdir=initialdir)
            else:
                root = tk.Tk()
                root.withdraw()
                try:
                    root.attributes('-topmost', True)
                except Exception:
                    pass
                folder = filedialog.askdirectory(title=title, initialdir=initialdir)
                root.destroy()
            if folder:  # 原生对话框成功返回了路径
                return folder
        except Exception:
            pass

        # 原生对话框失败或被取消，使用自定义对话框
        result = [None]

        dialog = tk.Toplevel(parent) if parent else tk.Tk()
        dialog.title(title)
        dialog.geometry("500x420")
        dialog.resizable(False, False)
        if parent:
            dialog.transient(parent)
            dialog.grab_set()
            dialog.update_idletasks()
            x = parent.winfo_x() + (parent.winfo_width() - 500) // 2
            y = parent.winfo_y() + (parent.winfo_height() - 420) // 2
            dialog.geometry(f"+{max(0,x)}+{max(0,y)}")
        else:
            dialog.update_idletasks()
            x = (dialog.winfo_screenwidth() - 500) // 2
            y = (dialog.winfo_screenheight() - 420) // 2
            dialog.geometry(f"+{x}+{y}")

        # 标题栏
        title_bar = tk.Frame(dialog, bg='#4A90D9', height=40)
        title_bar.pack(fill=X)
        title_bar.pack_propagate(False)
        tk.Label(title_bar, text=title, font=("TkDefaultFont", 12, "bold"),
                 bg='#4A90D9', fg='white').pack(side=LEFT, padx=15, pady=8)

        # 提示
        tk.Label(dialog, text="请输入目标文件夹路径，或点击「浏览」选择：",
                 font=("TkDefaultFont", 10), bg='#F5F6FA', fg='#2C3E50',
                 wraplength=460, justify=LEFT).pack(anchor=W, padx=20, pady='15 5')

        # 路径输入框
        input_frame = tk.Frame(dialog, bg='#F5F6FA')
        input_frame.pack(fill=X, padx=20, pady=5)

        path_var = tk.StringVar(value=initialdir or os.path.expanduser("~"))
        path_entry = tk.Entry(input_frame, textvariable=path_var, font=("TkDefaultFont", 11),
                              relief=FLAT, bg='#F0F2F5', fg='#2C3E50', insertbackground='#2C3E50')
        path_entry.pack(side=LEFT, fill=X, expand=True, ipady=6)
        path_entry.select_range(0, tk.END)
        path_entry.focus_set()

        # 浏览按钮（尝试原生对话框）
        def browse():
            try:
                bfolder = filedialog.askdirectory(title="选择文件夹", parent=dialog,
                                                  initialdir=path_var.get())
                if bfolder:
                    path_var.set(bfolder)
            except Exception:
                pass

        tk.Button(input_frame, text="浏览", font=("TkDefaultFont", 10),
                  bg='#4A90D9', fg='white', relief=FLAT, padx=12, cursor="hand2",
                  command=browse).pack(side=LEFT, padx='8 0')

        # 快捷路径
        quick_frame = tk.Frame(dialog, bg='#F5F6FA')
        quick_frame.pack(fill=X, padx=20, pady='10 0')
        tk.Label(quick_frame, text="快捷选择:", font=("TkDefaultFont", 9),
                 bg='#F5F6FA', fg='#7F8C8D').pack(side=LEFT)

        shortcuts = [
            ("主目录", os.path.expanduser("~")),
            ("桌面", _get_desktop_path()),
        ]
        # 添加常见挂载点
        for name, path in [("/data", "/data"), ("/home", "/home"), ("/tmp", "/tmp")]:
            if os.path.isdir(path):
                shortcuts.append((name, path))

        def set_path(p):
            path_var.set(p)

        for label, path in shortcuts:
            if path and os.path.isdir(path):
                tk.Button(quick_frame, text=label, font=("TkDefaultFont", 8),
                          bg='#E1E4E8', fg='#2C3E50', relief=FLAT, padx=8, pady=1,
                          cursor="hand2",
                          command=lambda p=path: set_path(p)).pack(side=LEFT, padx=2)

        # 当前路径状态
        status_var = tk.StringVar(value="")
        status_label = tk.Label(dialog, textvariable=status_var, font=("TkDefaultFont", 9),
                                bg='#F5F6FA', fg='#5CB85C', wraplength=460, justify=LEFT)
        status_label.pack(anchor=W, padx=20, pady='8 0')

        def check_path(*args):
            p = path_var.get().strip()
            if not p:
                status_var.set("")
            elif os.path.isdir(p):
                status_var.set(f"[OK] 路径有效: {p}")
            else:
                status_var.set(f"[!] 路径不存在，保存时将自动创建")

        path_var.trace_add('write', check_path)
        check_path()

        # 按钮区
        btn_frame = tk.Frame(dialog, bg='#F5F6FA')
        btn_frame.pack(fill=X, padx=20, pady=15)

        def confirm():
            p = path_var.get().strip()
            if p:
                result[0] = p
            dialog.destroy()

        def cancel():
            result[0] = None
            dialog.destroy()

        tk.Button(btn_frame, text="确定", font=("TkDefaultFont", 10, "bold"),
                  bg='#5CB85C', fg='white', relief=FLAT, padx=30, pady=5,
                  cursor="hand2", command=confirm).pack(side=RIGHT, padx=5)

        tk.Button(btn_frame, text="取消", font=("TkDefaultFont", 10),
                  bg='#95A5A6', fg='white', relief=FLAT, padx=30, pady=5,
                  cursor="hand2", command=cancel).pack(side=RIGHT, padx=5)

        # 回车确认
        dialog.bind('<Return>', lambda e: confirm())

        dialog.protocol("WM_DELETE_WINDOW", cancel)

        if not parent:
            dialog.mainloop()

        return result[0] if result[0] else ""
    except Exception:
        return ""


def _get_desktop_path() -> str:
    """快速获取桌面路径"""
    home = os.path.expanduser("~")
    xdg = os.environ.get("XDG_DESKTOP_DIR", "")
    if xdg and os.path.isdir(xdg):
        return xdg
    for name in ("Desktop", "桌面", "desktop"):
        p = os.path.join(home, name)
        if os.path.isdir(p):
            return p
    return home


def get_system_font() -> str:
    """获取系统可用中文字体（跨平台）"""
    if IS_WINDOWS:
        return "Microsoft YaHei UI"
    elif IS_LINUX:
        # 麒麟/Linux 常见中文字体
        for font in ("Noto Sans CJK SC", "WenQuanYi Micro Hei", "WenQuanYi Zen Hei",
                      "Droid Sans Fallback", "AR PL UMing CN", "Noto Sans SC"):
            try:
                import tkinter as tk
                root = tk.Tk()
                root.withdraw()
                families = root.tk.call('font', 'families')
                root.destroy()
                if font in families:
                    return font
            except Exception:
                pass
        return "TkDefaultFont"
    elif IS_MACOS:
        return "PingFang SC"
    return "TkDefaultFont"
